"""
PPT导出模块
将生成的教学文档导出为PPTX格式
"""

from pathlib import Path
from typing import Dict, Any, List, Optional
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from main.models import TeachingDocument, PPTContent, SlideKeyPoints
from main.config import get_template_config


class PPTExporter:
    """PPT导出器"""

    def __init__(self):
        """初始化导出器"""
        self.template_config = get_template_config()
        self.colors = self._init_colors()

    def _init_colors(self) -> Dict[str, Dict[str, RGBColor]]:
        """初始化配色方案"""
        return {
            'professional': {
                'primary': RGBColor(25, 55, 120),      # 深蓝
                'secondary': RGBColor(100, 130, 180),  # 灰蓝
                'accent': RGBColor(220, 100, 80),      # 橙红
                'text': RGBColor(40, 40, 40),          # 深灰
                'light': RGBColor(240, 240, 245)       # 浅蓝灰
            },
            'educational': {
                'primary': RGBColor(50, 120, 180),     # 蓝
                'secondary': RGBColor(100, 180, 80),   # 绿
                'accent': RGBColor(255, 140, 60),      # 橙
                'text': RGBColor(50, 50, 50),          # 暗灰
                'light': RGBColor(245, 250, 245)       # 浅绿白
            },
            'creative': {
                'primary': RGBColor(180, 50, 120),     # 紫红
                'secondary': RGBColor(80, 150, 200),   # 蓝
                'accent': RGBColor(255, 100, 100),     # 粉红
                'text': RGBColor(60, 60, 60),          # 灰
                'light': RGBColor(250, 240, 245)       # 浅粉白
            }
        }

    def export_to_pptx(self, document: TeachingDocument, output_path: str) -> str:
        """导出为PPTX格式

        Args:
            document: 教学文档
            output_path: 输出路径

        Returns:
            输出文件路径
        """

        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # 获取模板和配色
        template_name = document.ppt_layout.template_name if document.ppt_layout else 'professional'
        colors = self.colors.get(template_name, self.colors['professional'])

        # 添加标题幻灯片
        self._add_title_slide(prs, document, colors)

        # 添加内容幻灯片（带图片）
        if document.ppt_content and document.ppt_layout:
            self._add_content_slides_with_images(prs, document, colors)
        elif document.ppt_content:
            self._add_content_slides(prs, document.ppt_content, colors)

        # 保存文件
        prs.save(output_path)
        print(f"✅ PPT已导出: {output_path}")

        return output_path

    def _add_title_slide(self, prs: Presentation, document: TeachingDocument, colors: Dict[str, RGBColor]):
        """添加标题幻灯片"""

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors['primary']

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True

        title_para = title_frame.paragraphs[0]
        title_para.text = document.title
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

        # 添加副标题
        if document.ppt_content:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2), Inches(9), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = document.ppt_content.subtitle or "教学演示"
            subtitle_para.font.size = Pt(28)
            subtitle_para.font.color.rgb = colors['light']
            subtitle_para.alignment = PP_ALIGN.CENTER

        # 添加日期
        date_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.5), Inches(9), Inches(0.8)
        )
        date_frame = date_box.text_frame
        date_para = date_frame.paragraphs[0]
        date_para.text = datetime.now().strftime("%Y年%m月%d日")
        date_para.font.size = Pt(14)
        date_para.font.color.rgb = colors['light']
        date_para.alignment = PP_ALIGN.CENTER

    def _add_content_slides(self, prs: Presentation, ppt_content: PPTContent,
                           colors: Dict[str, RGBColor]):
        """添加内容幻灯片"""

        for slide_data in ppt_content.slides:
            self._add_single_slide(prs, slide_data, colors)

    def _add_single_slide(self, prs: Presentation, slide_data: SlideKeyPoints,
                         colors: Dict[str, RGBColor]):
        """添加单张幻灯片"""

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

        # 背景颜色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        # 添加顶部条纹
        header_shape = slide.shapes.add_shape(
            1,  # 矩形
            Inches(0), Inches(0), Inches(10), Inches(0.8)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = colors['primary']
        header_shape.line.color.rgb = colors['primary']

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True

        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.title
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.LEFT

        # 添加关键要点
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.2), Inches(8.4), Inches(5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        for i, key_point in enumerate(slide_data.key_points):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()

            p.text = f"• {key_point}"
            p.font.size = Pt(20)
            p.font.color.rgb = colors['text']
            p.level = 0
            p.space_before = Pt(6)
            p.space_after = Pt(6)

        # 添加底部说明（如果有）
        if slide_data.speaker_notes:
            footer_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(6.5), Inches(8.4), Inches(0.8)
            )
            footer_frame = footer_box.text_frame
            footer_frame.word_wrap = True

            footer_para = footer_frame.paragraphs[0]
            footer_para.text = f"💡 {slide_data.speaker_notes[:100]}..."
            footer_para.font.size = Pt(12)
            footer_para.font.italic = True
            footer_para.font.color.rgb = colors['secondary']

    def _add_content_slides_with_images(self, prs: Presentation, document: 'TeachingDocument',
                                       colors: Dict[str, RGBColor]):
        """添加包含图片的内容幻灯片"""
        from main.materials import IndexManagementService
        from main.config import get_storage_config
        
        material_service = IndexManagementService()
        storage_config = get_storage_config()
        materials_dir = Path(storage_config['materials_dir'])
        
        for idx, slide_data in enumerate(document.ppt_content.slides):
            # 找到对应的布局信息
            slide_layout = None
            if document.ppt_layout and idx < len(document.ppt_layout.slides):
                slide_layout = document.ppt_layout.slides[idx]
            
            # 获取图片ID
            image_ids = slide_layout.image_ids if slide_layout else []
            
            # 创建幻灯片
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
            
            # 背景颜色
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)
            
            # 添加顶部条纹
            header_shape = slide.shapes.add_shape(
                1,  # 矩形
                Inches(0), Inches(0), Inches(10), Inches(0.8)
            )
            header_shape.fill.solid()
            header_shape.fill.fore_color.rgb = colors['primary']
            header_shape.line.color.rgb = colors['primary']
            
            # 添加标题
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.15), Inches(9), Inches(0.6)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            
            title_para = title_frame.paragraphs[0]
            title_para.text = slide_data.title
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255)
            title_para.alignment = PP_ALIGN.LEFT
            
            # 如果有图片，调整布局
            has_image = bool(image_ids)
            
            if has_image:
                # 左侧文字，右侧图片布局
                content_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(1.2), Inches(4.5), Inches(5)
                )
                
                # 添加图片
                for img_id in image_ids[:1]:  # 只使用第一张图片
                    material = material_service.get_material_by_id(img_id)
                    if material:
                        image_path = materials_dir / material['filename']
                        if image_path.exists():
                            try:
                                # 在右侧添加图片
                                slide.shapes.add_picture(
                                    str(image_path),
                                    Inches(5.5), Inches(1.5),
                                    width=Inches(3.8)
                                )
                                print(f"✅ 为幻灯片 {idx} ({slide_data.title}) 添加图片: {material['filename']}")
                            except Exception as e:
                                print(f"⚠️  添加图片失败: {e}")
            else:
                # 无图片，全宽文字
                content_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(1.2), Inches(8.4), Inches(5)
                )
            
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            # 添加关键要点
            for i, key_point in enumerate(slide_data.key_points):
                if i == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()
                
                p.text = f"• {key_point}"
                p.font.size = Pt(20)
                p.font.color.rgb = colors['text']
                p.level = 0
                p.space_before = Pt(6)
                p.space_after = Pt(6)
            
            # 添加底部说明（如果有）
            if slide_data.speaker_notes:
                footer_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(6.5), Inches(8.4), Inches(0.8)
                )
                footer_frame = footer_box.text_frame
                footer_frame.word_wrap = True
                
                footer_para = footer_frame.paragraphs[0]
                footer_para.text = f"💡 {slide_data.speaker_notes[:100]}..."
                footer_para.font.size = Pt(12)
                footer_para.font.italic = True
                footer_para.font.color.rgb = colors['secondary']

    def export_to_markdown(self, document: TeachingDocument, output_path: str) -> str:
        """导出为Markdown格式"""

        md_content = self._generate_markdown(document)

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(md_content)

        print(f"✅ Markdown已导出: {output_path}")
        return output_path

    def _generate_markdown(self, document: TeachingDocument) -> str:
        """生成Markdown内容"""

        md = f"# {document.title}\n\n"

        # 添加教学意图
        if document.intent:
            md += "## 教学意图\n\n"
            md += f"**主题**: {document.intent.topic}\n\n"
            md += "**教学目标**:\n"
            for obj in document.intent.objectives:
                md += f"- {obj}\n"
            md += f"\n**受众层级**: {document.intent.audience_level.value}\n\n"

        # 添加教案
        if document.teaching_plan:
            md += "## 教案\n\n"
            md += f"### 导入\n{document.teaching_plan.introduction}\n\n"
            md += f"### 正文\n{document.teaching_plan.content}\n\n"
            md += "### 重点难点\n"
            for kp in document.teaching_plan.key_points:
                md += f"- {kp}\n"
            md += "\n"
            if document.teaching_plan.classroom_activities:
                md += f"### 课堂活动\n{document.teaching_plan.classroom_activities}\n\n"
            if document.teaching_plan.homework:
                md += f"### 课堂作业\n{document.teaching_plan.homework}\n\n"

        # 添加PPT大纲
        if document.ppt_content:
            md += "## PPT大纲\n\n"
            for i, slide in enumerate(document.ppt_content.slides, 1):
                md += f"### 幻灯片{i}: {slide.title}\n\n"
                md += "关键要点:\n"
                for kp in slide.key_points:
                    md += f"- {kp}\n"
                if slide.speaker_notes:
                    md += f"\n**说话稿**: {slide.speaker_notes}\n"
                md += "\n"

        return md


class ExportManager:
    """导出管理器"""

    def __init__(self):
        """初始化导出管理器"""
        self.ppt_exporter = PPTExporter()

    def export_document(self, document: TeachingDocument, format: str = "pptx",
                       output_path: Optional[str] = None) -> str:
        """导出文档

        Args:
            document: 教学文档
            format: 导出格式 (pptx/pdf/md)
            output_path: 输出路径（可选）

        Returns:
            导出文件路径
        """

        if not output_path:
            output_dir = Path('./data/outputs')
            output_dir.mkdir(parents=True, exist_ok=True)
            output_path = output_dir / f"{document.doc_id}.{format}"

        if format == 'pptx':
            return self.ppt_exporter.export_to_pptx(document, output_path)
        elif format == 'md':
            return self.ppt_exporter.export_to_markdown(document, output_path)
        elif format == 'pdf':
            # PDF导出需要额外的库，这里暂时使用Markdown作为替代
            return self.ppt_exporter.export_to_markdown(document, output_path.replace('.pdf', '.md'))
        else:
            raise ValueError(f"不支持的导出格式: {format}")

